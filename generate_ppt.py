"""Generate AI Governance Platform presentation following Ampera template."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import copy

# Ampera theme colors
C_MAGENTA   = RGBColor(0x95, 0x13, 0x5A)  # accent1 - primary brand
C_LPURPLE   = RGBColor(0xAC, 0x9B, 0xE9)  # accent4 - light purple
C_MPURPLE   = RGBColor(0x89, 0x71, 0xE1)  # accent5 - medium purple
C_DARK      = RGBColor(0x45, 0x45, 0x51)  # dk2 - dark gray
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LAVENDER  = RGBColor(0xCF, 0xC6, 0xF3)  # accent6 - very light purple
C_LIGHT_BG  = RGBColor(0xF5, 0xF0, 0xFF)  # near-white lavender for boxes

FONT = "Open Sans"

prs = Presentation(
    "/home/ampara/Documents/aigovernance_backend/Ampera Technologies - PPT 2.pptx"
)

layouts = {l.name: l for l in prs.slide_master.slide_layouts}
BLANK_L       = layouts["2_Blank"]         # AI-background + top-right logo, no bottom triangles
CONTENT_L     = layouts["Title and Content"]  # left triangles + top-right logo
TWO_CONTENT_L = layouts["Two Content"]     # two columns (no triangles/logo — added manually)


# ── helpers ──────────────────────────────────────────────────────────────────

def set_font(run, size_pt, bold=False, color=None, name=FONT):
    run.font.name  = name
    run.font.size  = Pt(size_pt)
    run.font.bold  = bold
    if color:
        run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size_pt, bold=False, color=C_DARK,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = wrap
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold=bold, color=color)
    run.font.italic = italic
    return tb


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width_pt=0):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt) if line_width_pt else Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rect_with_text(slide, x, y, w, h, title, body_lines,
                        fill=C_LIGHT_BG, border=C_LPURPLE,
                        title_size=10, body_size=8.5,
                        title_color=C_MAGENTA, body_color=C_DARK):
    rect = add_rect(slide, x, y, w, h, fill_color=fill,
                    line_color=border, line_width_pt=1.0)
    tf = rect.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn as _qn
    tf.margin_left  = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top   = Inches(0.05)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    set_font(run, title_size, bold=True, color=title_color)

    for line in body_lines:
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsmap
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = line
        set_font(run2, body_size, color=body_color)
    return rect


def add_arrow_right(slide, x, y, length=0.35, shaft_h=0.05, color=C_MPURPLE):
    """Simple right-pointing arrow as a filled shape."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    # Use an auto-shape: right arrow (shape type 13 = right arrow)
    shape = slide.shapes.add_shape(
        13,  # right arrow
        Inches(x), Inches(y - 0.12), Inches(length), Inches(0.25)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def bullet_para(tf, text, level=0, size_pt=14, color=C_DARK, bold=False):
    p = tf.add_paragraph()
    p.level = level
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold=bold, color=color)
    return p


def set_title(slide, text, size_pt=28, color=C_MAGENTA):
    ph = slide.placeholders[0]
    ph.text = ""
    tf = ph.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold=True, color=color)


def inject_side_triangles(slide):
    """Inject the left-side triangle decoration into any slide that lacks it."""
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    max_id = max((sh.shape_id for sh in slide.shapes), default=0)
    spTree = slide.element.find(f'.//{{{NS_P}}}spTree')

    # Triangle 8 — magenta (accent1), full-height, rotated 180° + flipH
    tri8 = etree.fromstring(
        f'<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">'
        f'<p:nvSpPr><p:cNvPr id="{max_id+1}" name="Right Triangle 8"/>'
        f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm rot="10800000" flipH="1">'
        f'<a:off x="1" y="-4"/><a:ext cx="838200" cy="6858003"/></a:xfrm>'
        f'<a:prstGeom prst="rtTriangle"><a:avLst/></a:prstGeom>'
        f'<a:ln><a:noFill/></a:ln></p:spPr>'
        f'<p:style>'
        f'<a:lnRef idx="2"><a:schemeClr val="accent1"><a:shade val="15000"/></a:schemeClr></a:lnRef>'
        f'<a:fillRef idx="1"><a:schemeClr val="accent1"/></a:fillRef>'
        f'<a:effectRef idx="0"><a:schemeClr val="accent1"/></a:effectRef>'
        f'<a:fontRef idx="minor"><a:schemeClr val="lt1"/></a:fontRef>'
        f'</p:style>'
        f'<p:txBody><a:bodyPr rtlCol="0" anchor="ctr"/><a:lstStyle/><a:p/></p:txBody>'
        f'</p:sp>'
    )

    # Triangle 9 — light purple (accent4), from below title area down to bottom
    tri9 = etree.fromstring(
        f'<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">'
        f'<p:nvSpPr><p:cNvPr id="{max_id+2}" name="Right Triangle 9"/>'
        f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="-1" y="1066800"/>'
        f'<a:ext cx="838201" cy="5796024"/></a:xfrm>'
        f'<a:prstGeom prst="rtTriangle"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:schemeClr val="accent4"/></a:solidFill>'
        f'<a:ln><a:noFill/></a:ln></p:spPr>'
        f'<p:txBody><a:bodyPr rtlCol="0" anchor="ctr"/><a:lstStyle/><a:p/></p:txBody>'
        f'</p:sp>'
    )

    # Insert at the front of the draw order so triangles sit behind all other content
    spTree.insert(2, tri8)
    spTree.insert(3, tri9)


def add_logo_topright(slide):
    """Add the Ampera logo to the top-right corner (matching layout position)."""
    slide.shapes.add_picture(
        '/tmp/image1.jpeg',
        Inches(11.49), Inches(0.02), Inches(1.85), Inches(0.49)
    )


# ── delete the initial blank slide (index 0) ─────────────────────────────────

def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    sldId = xml_slides[index]
    rId = sldId.get(qn("r:id"))
    xml_slides.remove(sldId)
    try:
        del prs.part.related_parts[rId]
    except Exception:
        pass
    try:
        rel_elem = prs.part.rels[rId]._element
        rel_elem.getparent().remove(rel_elem)
        del prs.part._rels[rId]
    except Exception:
        pass


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Overview  (Blank layout → AI background + Ampera logo centered)
# ═══════════════════════════════════════════════════════════════════════════

s1 = prs.slides.add_slide(BLANK_L)
inject_side_triangles(s1)   # add left-side bar (same as content slides)

# Main title — centered, large, white
add_textbox(s1, 1.5, 1.4, 10.33, 1.5,
            "AI Governance Platform",
            size_pt=44, bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER)

# Subtitle
add_textbox(s1, 1.5, 3.1, 10.33, 0.8,
            "Centralized control layer for enterprise AI operations",
            size_pt=20, color=C_LAVENDER,
            align=PP_ALIGN.CENTER)

# Divider-style tagline
add_textbox(s1, 1.5, 4.05, 10.33, 0.5,
            "Secure  ·  Scalable  ·  Compliant",
            size_pt=14, color=C_LPURPLE,
            align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Architecture Diagram  (Title and Content layout)
# ═══════════════════════════════════════════════════════════════════════════

s2 = prs.slides.add_slide(CONTENT_L)
set_title(s2, "System Architecture", size_pt=26)

# Clear the content placeholder
s2.placeholders[1].text = ""

# ── Row of boxes: Client → Gateway → Provider ─────────────────────────────
# Available horizontal space (after left triangle): x≈1.1 to 13.0"
# Keep all within: x=1.1", total width=11.5", heights centred around y=3.0"

# Client Teams box
add_rect_with_text(s2, 1.15, 2.35, 1.95, 1.8,
                   "Client\nTeams",
                   ["Enterprise apps,", "internal tools &", "AI consumers"],
                   fill=C_LIGHT_BG, border=C_LPURPLE,
                   title_size=10, body_size=8)

# Arrow 1
add_arrow_right(s2, 3.18, 3.28, length=0.38, color=C_MAGENTA)

# Central Gateway box (taller)
add_rect_with_text(s2, 3.65, 1.9, 5.2, 2.7,
                   "AI Governance Platform",
                   [
                       "────────────────────",
                       "Policy & Access Control",
                       "PII Protection",
                       "Budget & Rate Limiting",
                       "Audit Logging",
                   ],
                   fill=RGBColor(0xF8, 0xF0, 0xFF),
                   border=C_MAGENTA,
                   title_size=11, body_size=8.5,
                   title_color=C_MAGENTA, body_color=C_DARK)

# Arrow 2
add_arrow_right(s2, 8.93, 3.28, length=0.38, color=C_MAGENTA)

# AI Provider box
add_rect_with_text(s2, 9.40, 2.35, 2.2, 1.8,
                   "AI Provider\n(Azure OpenAI)",
                   ["Processes", "approved AI", "requests"],
                   fill=C_LIGHT_BG, border=C_LPURPLE,
                   title_size=10, body_size=8)

# ── Down arrow + Analytics box ─────────────────────────────────────────────
# Down arrow from Gateway box
down_arrow = s2.shapes.add_shape(
    14,  # down arrow
    Inches(5.8), Inches(4.62), Inches(0.25), Inches(0.35)
)
down_arrow.fill.solid()
down_arrow.fill.fore_color.rgb = C_MPURPLE
down_arrow.line.fill.background()

add_rect_with_text(s2, 3.65, 4.97, 5.2, 1.0,
                   "Dashboard & Analytics",
                   ["Real-time visibility into usage, cost, and compliance"],
                   fill=RGBColor(0xEC, 0xE8, 0xFF), border=C_MPURPLE,
                   title_size=10, body_size=8.5,
                   title_color=C_MPURPLE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Proxy Setup  (Title and Content)
# ═══════════════════════════════════════════════════════════════════════════

s3 = prs.slides.add_slide(CONTENT_L)
set_title(s3, "Intelligent Proxy Gateway", size_pt=26)

ph3 = s3.placeholders[1]
ph3.text = ""
tf3 = ph3.text_frame
tf3.word_wrap = True

bullets3 = [
    ("Single entry point for all AI API traffic",
     "Every request from every team passes through one controlled gateway before reaching the AI provider."),
    ("Real-time policy enforcement",
     "Identity is verified and access rules are applied instantly — models, token limits, and data policies."),
    ("Automatic PII masking",
     "Sensitive information is detected and redacted before any data leaves the organization's boundary."),
    ("Transparent to end users",
     "Teams interact with their existing workflows; the governance layer operates silently in the background."),
]

first = True
for heading, detail in bullets3:
    if first:
        p = tf3.paragraphs[0]
        first = False
    else:
        p = tf3.add_paragraph()
    p.level = 0
    run = p.add_run()
    run.text = heading
    set_font(run, 14, bold=True, color=C_MAGENTA)

    p2 = tf3.add_paragraph()
    p2.level = 1
    run2 = p2.add_run()
    run2.text = detail
    set_font(run2, 11.5, color=C_DARK)

    # spacer
    tf3.add_paragraph()


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Cost & Budget Management  (Title and Content)
# ═══════════════════════════════════════════════════════════════════════════

s4 = prs.slides.add_slide(CONTENT_L)
set_title(s4, "Cost & Budget Management", size_pt=26)

ph4 = s4.placeholders[1]
ph4.text = ""
tf4 = ph4.text_frame
tf4.word_wrap = True

bullets4 = [
    ("Real-time spend tracking",
     "Costs are calculated instantly for every request, broken down by team, project, and AI model."),
    ("Enforced monthly budgets",
     "Administrators set spending ceilings per team or project; requests are blocked automatically when limits are hit."),
    ("Proactive alerts",
     "Configurable threshold notifications are sent before limits are reached, preventing unexpected disruptions."),
    ("Transparent cost reporting",
     "Executives and finance teams get clear, aggregated usage reports with full traceability to individual teams."),
]

first = True
for heading, detail in bullets4:
    if first:
        p = tf4.paragraphs[0]
        first = False
    else:
        p = tf4.add_paragraph()
    p.level = 0
    run = p.add_run()
    run.text = heading
    set_font(run, 14, bold=True, color=C_MAGENTA)

    p2 = tf4.add_paragraph()
    p2.level = 1
    run2 = p2.add_run()
    run2.text = detail
    set_font(run2, 11.5, color=C_DARK)

    tf4.add_paragraph()


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Platform Modules  (Two Content layout)
# ═══════════════════════════════════════════════════════════════════════════

s5 = prs.slides.add_slide(TWO_CONTENT_L)
set_title(s5, "Platform Modules", size_pt=26)

# Left column placeholder (idx=1)
left_modules = [
    ("Governance Rules",
     "Controls which AI models and capabilities each team is permitted to access."),
    ("PII Protection",
     "Automatically detects and masks sensitive personal data before it reaches the AI provider."),
    ("Rate Limiting",
     "Prevents any single team from consuming a disproportionate share of shared resources."),
]

# Right column placeholder (idx=2)
right_modules = [
    ("Budget Control",
     "Enforces monthly spend limits and blocks requests when allocations are exceeded."),
    ("Security & Alerts",
     "Monitors for anomalous usage patterns and delivers notifications to administrators."),
    ("Audit Logging",
     "Records every request and governance decision, providing a complete compliance trail."),
]

def fill_module_column(slide, ph_idx, modules):
    ph = slide.placeholders[ph_idx]
    ph.text = ""
    tf = ph.text_frame
    tf.word_wrap = True
    first = True
    for mod_name, mod_desc in modules:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = mod_name
        set_font(run, 13, bold=True, color=C_MAGENTA)

        p2 = tf.add_paragraph()
        p2.level = 1
        run2 = p2.add_run()
        run2.text = mod_desc
        set_font(run2, 10.5, color=C_DARK)

        tf.add_paragraph()

fill_module_column(s5, 1, left_modules)
fill_module_column(s5, 2, right_modules)
inject_side_triangles(s5)   # Two Content layout has no triangles — add them
add_logo_topright(s5)       # Two Content layout has no logo — add it


# ═══════════════════════════════════════════════════════════════════════════
# Delete the original blank slide (now at index 0)
# ═══════════════════════════════════════════════════════════════════════════

delete_slide(prs, 0)

# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════

out_path = "/home/ampara/Documents/aigovernance_backend/Ampera AI Governance - Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
